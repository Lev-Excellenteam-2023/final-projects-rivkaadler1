from io import BytesIO
from pptx import Presentation


def parse_from_binary_file_to_pptx(path: str):
    """
     Parses a PowerPoint presentation from a binary file and returns the corresponding presentation object.

     Parameters:
         path (str): The file path of the binary PowerPoint presentation to be parsed.

     Return Value:
         presentation (Presentation): A presentation object representing the parsed PowerPoint presentation.
    """
    with open(path, "rb") as file:
        pptx_data = file.read()
    presentation = Presentation(BytesIO(pptx_data))
    return presentation


def parse_content_to_string_dict(presentation: Presentation):
    """
    Parses the content of a PowerPoint file and returns a dictionary with slide numbers as keys and
    slide content (title, content, notes, and slide number) as string values.

    :param presentation:
    :return: A dictionary with slide numbers as keys and slide
    content (title, content, notes, and slide number) as string values.
    """
    parsed_content = {}

    for i, slide in enumerate(presentation.slides, start=1):
        title = slide.shapes.title.text.strip() if slide.shapes.title else ""
        content = ''.join(run.text.strip() for shape in slide.shapes if shape.has_text_frame
                          for paragraph in shape.text_frame.paragraphs for run in paragraph.runs)

        notes_slide = slide.notes_slide
        notes = ''.join(run.text.strip() for shape in notes_slide.shapes if shape.has_text_frame
                        for paragraph in shape.text_frame.paragraphs for run in paragraph.runs)

        if title or content or notes:
            slide_content = f"title: {title}\ncontent: {content}\nnotes: {notes}\nslide number: {i}"
            parsed_content[i] = slide_content

    return parsed_content
