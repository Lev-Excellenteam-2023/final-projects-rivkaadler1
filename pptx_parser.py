from pptx import Presentation


def parse_content_to_string_dict(pptx_path: str):
    """
    Parses the content of a PowerPoint file and returns a dictionary with slide numbers as keys and
    slide content (title, content, notes, and slide number) as string values.

    :param pptx_path: The path to the PowerPoint file.
    :return: A dictionary with slide numbers as keys and slide
    content (title, content, notes, and slide number) as string values.
    """
    presentation = Presentation(pptx_path)
    parsed_content = {}

    for i, slide in enumerate(presentation.slides, start=1):
        title = slide.shapes.title.text.strip() if slide.shapes.title else ""
        content = ''.join(run.text.strip() for shape in slide.shapes if shape.has_text_frame
                          for paragraph in shape.text_frame.paragraphs for run in paragraph.runs)

        notes_slide = slide.notes_slide
        notes = ''.join(run.text.strip() for shape in notes_slide.shapes if shape.has_text_frame
                        for paragraph in shape.text_frame.paragraphs for run in paragraph.runs)

        if title or content or notes:
            slide_content = f"title: {title}\ncontent: {content}\nnotes: {notes}\nslide number: {i}"
            parsed_content[i] = slide_content

    return parsed_content

